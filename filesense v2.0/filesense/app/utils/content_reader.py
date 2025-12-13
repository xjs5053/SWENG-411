"""
Content Reader - Extract text content from various file types
"""
import os
from pathlib import Path
from typing import Optional, Tuple
import chardet


class ContentReader:
    """Extract text content from various file formats"""
    
    # Maximum file size to read (in bytes) - 10MB
    MAX_FILE_SIZE = 10 * 1024 * 1024
    
    # Maximum characters to return
    MAX_CHARS = 50000
    
    @staticmethod
    def read_file(file_path: str, max_chars: int = None) -> Tuple[str, str]:
        """
        Read file content and return (content, error_message)
        Returns empty string and error message if failed
        """
        if max_chars is None:
            max_chars = ContentReader.MAX_CHARS
        
        if not os.path.exists(file_path):
            return "", "File not found"
        
        file_size = os.path.getsize(file_path)
        if file_size > ContentReader.MAX_FILE_SIZE:
            return "", f"File too large ({file_size / (1024*1024):.1f} MB)"
        
        ext = Path(file_path).suffix.lower()
        
        try:
            # Text files
            if ext in ['.txt', '.md', '.log', '.csv', '.json', '.xml', '.yaml', '.yml',
                      '.ini', '.cfg', '.conf', '.sh', '.bash', '.sql', '.html', '.css',
                      '.js', '.ts', '.jsx', '.tsx', '.py', '.java', '.cs', '.cpp', '.c',
                      '.h', '.rb', '.go', '.rs', '.php', '.swift', '.kt']:
                return ContentReader._read_text_file(file_path, max_chars)
            
            # PDF files
            elif ext == '.pdf':
                return ContentReader._read_pdf(file_path, max_chars)
            
            # Word documents
            elif ext in ['.docx', '.doc']:
                return ContentReader._read_docx(file_path, max_chars)
            
            # Excel files
            elif ext in ['.xlsx', '.xls']:
                return ContentReader._read_xlsx(file_path, max_chars)
            
            # PowerPoint files
            elif ext in ['.pptx', '.ppt']:
                return ContentReader._read_pptx(file_path, max_chars)
            
            else:
                # Try as text
                return ContentReader._read_text_file(file_path, max_chars)
                
        except Exception as e:
            return "", f"Error reading file: {str(e)}"
    
    @staticmethod
    def _read_text_file(file_path: str, max_chars: int) -> Tuple[str, str]:
        """Read plain text file with encoding detection"""
        try:
            # Detect encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read(min(10000, os.path.getsize(file_path)))
                result = chardet.detect(raw_data)
                encoding = result.get('encoding', 'utf-8') or 'utf-8'
            
            # Read file with detected encoding
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                content = f.read(max_chars)
            
            return content, ""
            
        except Exception as e:
            # Fallback to utf-8 with error handling
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(max_chars)
                return content, ""
            except Exception as e2:
                return "", f"Could not read text file: {str(e2)}"
    
    @staticmethod
    def _read_pdf(file_path: str, max_chars: int) -> Tuple[str, str]:
        """Read PDF file"""
        try:
            from PyPDF2 import PdfReader
            
            reader = PdfReader(file_path)
            text_parts = []
            total_chars = 0
            
            for page in reader.pages:
                page_text = page.extract_text() or ""
                text_parts.append(page_text)
                total_chars += len(page_text)
                
                if total_chars >= max_chars:
                    break
            
            content = "\n\n".join(text_parts)[:max_chars]
            return content, ""
            
        except ImportError:
            return "", "PyPDF2 not installed. Run: pip install PyPDF2"
        except Exception as e:
            return "", f"Error reading PDF: {str(e)}"
    
    @staticmethod
    def _read_docx(file_path: str, max_chars: int) -> Tuple[str, str]:
        """Read Word document"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = []
            total_chars = 0
            
            for para in doc.paragraphs:
                text_parts.append(para.text)
                total_chars += len(para.text)
                
                if total_chars >= max_chars:
                    break
            
            # Also extract from tables
            if total_chars < max_chars:
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        text_parts.append(row_text)
                        total_chars += len(row_text)
                        
                        if total_chars >= max_chars:
                            break
            
            content = "\n".join(text_parts)[:max_chars]
            return content, ""
            
        except ImportError:
            return "", "python-docx not installed. Run: pip install python-docx"
        except Exception as e:
            return "", f"Error reading Word document: {str(e)}"
    
    @staticmethod
    def _read_xlsx(file_path: str, max_chars: int) -> Tuple[str, str]:
        """Read Excel spreadsheet"""
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            total_chars = 0
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows(max_row=100):  # Limit rows
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    
                    if row_values:
                        row_text = " | ".join(row_values)
                        text_parts.append(row_text)
                        total_chars += len(row_text)
                    
                    if total_chars >= max_chars:
                        break
                
                if total_chars >= max_chars:
                    break
            
            wb.close()
            content = "\n".join(text_parts)[:max_chars]
            return content, ""
            
        except ImportError:
            return "", "openpyxl not installed. Run: pip install openpyxl"
        except Exception as e:
            return "", f"Error reading Excel file: {str(e)}"
    
    @staticmethod
    def _read_pptx(file_path: str, max_chars: int) -> Tuple[str, str]:
        """Read PowerPoint presentation"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            total_chars = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"=== Slide {slide_num} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
                        total_chars += len(shape.text)
                
                if total_chars >= max_chars:
                    break
            
            content = "\n".join(text_parts)[:max_chars]
            return content, ""
            
        except ImportError:
            return "", "python-pptx not installed. Run: pip install python-pptx"
        except Exception as e:
            return "", f"Error reading PowerPoint: {str(e)}"
    
    @staticmethod
    def get_preview(file_path: str, lines: int = 50) -> str:
        """Get a preview of file content (first N lines)"""
        content, error = ContentReader.read_file(file_path, max_chars=10000)
        
        if error:
            return f"[Preview unavailable: {error}]"
        
        if not content:
            return "[File is empty]"
        
        content_lines = content.split('\n')[:lines]
        return '\n'.join(content_lines)
    
    @staticmethod
    def can_read(file_path: str) -> bool:
        """Check if content reader can handle this file type"""
        supported = [
            '.txt', '.md', '.log', '.csv', '.json', '.xml', '.yaml', '.yml',
            '.ini', '.cfg', '.conf', '.sh', '.bash', '.sql', '.html', '.css',
            '.js', '.ts', '.jsx', '.tsx', '.py', '.java', '.cs', '.cpp', '.c',
            '.h', '.rb', '.go', '.rs', '.php', '.swift', '.kt',
            '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'
        ]
        ext = Path(file_path).suffix.lower()
        return ext in supported
    
    @staticmethod
    def get_word_count(content: str) -> int:
        """Count words in content"""
        return len(content.split())
    
    @staticmethod
    def get_line_count(content: str) -> int:
        """Count lines in content"""
        return len(content.split('\n'))
